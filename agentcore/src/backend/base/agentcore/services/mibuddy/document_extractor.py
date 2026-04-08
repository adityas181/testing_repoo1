"""Document text extraction service.

Extracts text content from 22+ file types for document Q&A.
Ported from MiBuddy's extract_text_from_file() with agentcore storage integration.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path

from loguru import logger

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}

SUPPORTED_DOC_EXTENSIONS = {
    # Documents
    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".xlsm",
    # Text / Markup
    ".txt", ".md", ".csv", ".html", ".tex",
    # Code
    ".py", ".js", ".ts", ".java", ".cpp", ".c", ".cs", ".go",
    ".rb", ".php", ".sh", ".css", ".json",
}


async def read_file_bytes(file_path: str) -> bytes:
    """Read file bytes from storage.

    Tries multiple storage locations:
    1. Main storage (files uploaded via /api/files)
    2. MiBuddy dedicated container (uploads folder)

    file_path format: "{user_id}/{filename}".
    """
    # Try MiBuddy container first (with full path)
    try:
        from agentcore.services.mibuddy.docqa_storage import get_file_by_path
        data = await get_file_by_path(file_path)
        logger.info(f"[DocExtractor] Read {len(data)} bytes from MiBuddy container: {file_path}")
        return data
    except Exception as e:
        logger.debug(f"[DocExtractor] Not in MiBuddy container ({e}), trying main storage: {file_path}")

    # Fallback: try main storage
    parts = file_path.replace("\\", "/").split("/", 1)
    if len(parts) == 2:
        agent_id, file_name = parts
    else:
        agent_id, file_name = "", file_path

    try:
        from agentcore.services.deps import get_storage_service
        storage = get_storage_service()
        return await storage.get_file(agent_id=agent_id, file_name=file_name)
    except Exception:
        raise FileNotFoundError(f"File not found in any storage: {file_path}")


def extract_text_from_bytes(file_bytes: bytes, file_ext: str) -> str:
    """Extract text content from file bytes based on extension.

    Args:
        file_bytes: Raw file content.
        file_ext: File extension including dot (e.g. ".pdf").

    Returns:
        Extracted text content.
    """
    file_ext = file_ext.lower()

    try:
        # --- Text / Code files ---
        if file_ext in {
            ".txt", ".md", ".csv", ".c", ".cpp", ".cs", ".css", ".go",
            ".java", ".js", ".json", ".php", ".py", ".rb", ".sh",
            ".tex", ".ts", ".html",
        }:
            text = file_bytes.decode("utf-8", errors="ignore")
            if file_ext == ".html":
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(text, "html.parser")
                text = soup.get_text(separator="\n")
            return text

        # --- PDF ---
        if file_ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text

        # --- Word (DOCX) ---
        if file_ext == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        # --- PowerPoint (PPTX) ---
        if file_ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in paragraph.runs)
                            if line.strip():
                                text += line + "\n"
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text += cell.text + "\n"
            return text

        # --- Excel (XLS, XLSX, XLSM) ---
        if file_ext in {".xls", ".xlsx", ".xlsm"}:
            import pandas as pd
            file_stream = io.BytesIO(file_bytes)
            engine = "xlrd" if file_ext == ".xls" else "openpyxl"
            try:
                xls = pd.ExcelFile(file_stream, engine=engine)
            except Exception:
                file_stream.seek(0)
                xls = pd.ExcelFile(file_stream)

            parts: list[str] = []
            for sheet_idx, sheet_name in enumerate(xls.sheet_names, start=1):
                try:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    parts.append(f"\n--- Sheet {sheet_idx}: {sheet_name} ---\n")
                    parts.append(df.to_csv(index=False))
                except Exception as e:
                    logger.warning(f"Skipping Excel sheet '{sheet_name}': {e}")
            return "\n".join(parts).strip()

        return f"[Unsupported file type: {file_ext}]"

    except Exception as e:
        logger.error(f"Error extracting text from {file_ext}: {e}")
        return f"[ERROR] Unable to process file ({file_ext}). Reason: {str(e)}"


async def extract_text(file_path: str) -> str:
    """Read a file from storage and extract its text content.

    Args:
        file_path: Storage-relative path (e.g. "{user_id}/{filename}").

    Returns:
        Extracted text content.
    """
    ext = Path(file_path).suffix.lower()

    if ext in IMAGE_EXTENSIONS:
        return f"[Image file: {Path(file_path).name}]"

    if ext not in SUPPORTED_DOC_EXTENSIONS:
        return f"[Unsupported file type: {ext}]"

    try:
        file_bytes = await read_file_bytes(file_path)
        text = extract_text_from_bytes(file_bytes, ext)
        if not text.strip():
            return f"[No text content extracted from: {Path(file_path).name}]"
        return text
    except Exception as e:
        logger.error(f"Failed to extract text from {file_path}: {e}")
        return f"[ERROR] Failed to read file: {str(e)}"
