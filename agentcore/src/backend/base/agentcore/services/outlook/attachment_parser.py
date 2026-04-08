"""Parse email attachments to text — reuses project's existing document libraries.

Supported formats: txt, csv, pdf, docx, xlsx, pptx.
All other types return a clean error (no exception).
"""
from __future__ import annotations

import base64
import io
from typing import Any

from loguru import logger

# Max attachment size: 10 MB (base64 decoded)
MAX_ATTACHMENT_BYTES = 10 * 1024 * 1024
MAX_ATTACHMENTS = 20
SUPPORTED_EXTENSIONS = {".txt", ".csv", ".pdf", ".docx", ".xlsx", ".pptx"}


def parse_attachment(name: str, content_bytes_b64: str) -> dict[str, Any]:
    """Parse a single Graph API attachment to text.

    Args:
        name: Filename (e.g. "report.pdf")
        content_bytes_b64: Base64-encoded file content from Graph API

    Returns:
        {"filename": str, "text": str | None, "error": str | None, "size_bytes": int}
    """
    ext = ("." + name.rsplit(".", 1)[-1]).lower() if "." in name else ""
    result: dict[str, Any] = {"filename": name, "text": None, "error": None, "size_bytes": 0}

    try:
        raw = base64.b64decode(content_bytes_b64)
        result["size_bytes"] = len(raw)
    except Exception:
        result["error"] = "Failed to decode base64 content"
        return result

    if len(raw) > MAX_ATTACHMENT_BYTES:
        result["error"] = f"Attachment too large ({len(raw)} bytes, max {MAX_ATTACHMENT_BYTES})"
        return result

    if ext not in SUPPORTED_EXTENSIONS:
        result["error"] = f"Unsupported file type: {ext}"
        return result

    try:
        if ext == ".txt":
            result["text"] = raw.decode("utf-8", errors="replace")

        elif ext == ".csv":
            result["text"] = raw.decode("utf-8", errors="replace")

        elif ext == ".pdf":
            import pymupdf

            doc = pymupdf.open(stream=raw, filetype="pdf")
            result["text"] = "\n".join(page.get_text() for page in doc)
            doc.close()

        elif ext == ".docx":
            from docx import Document

            doc = Document(io.BytesIO(raw))
            result["text"] = "\n".join(p.text for p in doc.paragraphs)

        elif ext == ".xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    lines.append("\t".join(str(c) if c is not None else "" for c in row))
            result["text"] = "\n".join(lines)
            wb.close()

        elif ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(io.BytesIO(raw))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text)
            result["text"] = "\n".join(texts)

    except Exception as e:
        logger.warning("Failed to parse attachment {}: {}", name, e)
        result["error"] = f"Parse error: {e!s}"

    return result


def parse_attachments(attachments: list[dict]) -> list[dict[str, Any]]:
    """Parse a list of Graph API attachments (bounded by MAX_ATTACHMENTS).

    Args:
        attachments: List from Graph API /messages/{id}/attachments response

    Returns:
        List of parse results, one per attachment
    """
    results = []
    for att in attachments[:MAX_ATTACHMENTS]:
        # Graph API file attachments have @odata.type "#microsoft.graph.fileAttachment"
        if att.get("@odata.type") != "#microsoft.graph.fileAttachment":
            continue
        name = att.get("name", "unknown")
        content = att.get("contentBytes", "")
        results.append(parse_attachment(name, content))
    return results
