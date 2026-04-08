from __future__ import annotations

import base64
import mimetypes
import os
import traceback
from pathlib import Path

from loguru import logger

from agentcore.custom.custom_node.node import Node
from agentcore.io import BoolInput, HandleInput, IntInput, Output
from agentcore.schema.data import Data


class OCRExtractorNode(Node):
    """Extract text from documents using a connected multimodal LLM for OCR."""

    display_name: str = "Document OCR Extractor"
    description: str = (
        "Extract text from files using a connected multimodal Language Model. "
        "Supports PDF (native + scanned), images, DOCX, PPTX, XLSX, CSV, TXT."
    )
    name = "DocumentOCRExtractor"
    icon = "FileText"

    inputs = [
        HandleInput(
            name="llm",
            display_name="Language Model",
            input_types=["LanguageModel"],
            info="Multimodal LLM for vision-based OCR (e.g. GPT-4o, Gemini Flash, Claude). "
                 "Only used for scanned PDFs and images; native text is extracted without LLM.",
        ),
        HandleInput(
            name="file_paths",
            display_name="File Paths",
            input_types=["Data", "Message"],
            info="File paths from Knowledge Base component.",
            is_list=True,
        ),
        IntInput(name="min_native_text_length", display_name="Min Native Text Length", value=50, advanced=True),
        IntInput(name="ocr_dpi", display_name="OCR DPI", value=300, advanced=True),
        BoolInput(name="extract_tables", display_name="Extract Tables", value=True, advanced=True),
    ]

    outputs = [
        Output(
            display_name="Extracted Documents",
            name="documents",
            method="extract_documents",
            output_types=["Data"],
            is_list=True,
        ),
    ]

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp",
        ".docx", ".pptx", ".xlsx", ".csv", ".txt",
    }
    IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"}

    # ══════════════════════════════════════════════════════════
    #  MAIN ENTRY POINT
    # ══════════════════════════════════════════════════════════

    def extract_documents(self) -> list[Data]:
        # ── Resolve file paths ──
        paths = self._resolve_paths()
        logger.info(f"[OCR] Resolved {len(paths)} path(s)")

        if not paths:
            self.status = "No files found. Check Knowledge Base connection."
            return [Data(text="No files found to process.", data={"error": True})]

        # ── Extract text from each file, then merge pages into one document per file ──
        merged_docs: list[Data] = []
        errors: list[str] = []

        for path in paths:
            try:
                logger.info(f"[OCR] Extracting: {path}")
                page_docs = self._extract_file(path)

                if not page_docs:
                    continue

                # Merge all pages into a single document for proper chunking
                full_text = "\n\n".join(d.text for d in page_docs if d.text)
                if full_text.strip():
                    merged_docs.append(Data(text=full_text, data={
                        "source_file": path.name,
                        "file_path": str(path),
                        "total_pages": len(page_docs),
                        "file_type": path.suffix.lstrip(".").lower(),
                        "extraction_method": "ocr",
                    }))
                    logger.info(f"[OCR] Merged {len(page_docs)} page(s) into 1 document ({len(full_text)} chars)")
            except Exception as e:
                err_msg = f"{path.name}: {type(e).__name__}: {e}"
                logger.error(f"[OCR] Error extracting {path}: {err_msg}")
                logger.error(traceback.format_exc())
                errors.append(err_msg)

        status = f"Extracted {len(merged_docs)} document(s) from {len(paths)} file(s)"
        if errors:
            status += f" | Errors: {'; '.join(errors[:3])}"
        self.status = status

        if not merged_docs:
            return [Data(
                text=f"Extraction returned 0 documents. Errors: {'; '.join(errors)}",
                data={"error": True, "paths": [str(p) for p in paths]},
            )]

        return merged_docs

    # ══════════════════════════════════════════════════════════
    #  PATH RESOLUTION
    # ══════════════════════════════════════════════════════════

    def _resolve_paths(self) -> list[Path]:
        paths: list[Path] = []
        if not self.file_paths:
            return paths

        items = self.file_paths if isinstance(self.file_paths, list) else [self.file_paths]

        for item in items:
            if item is None:
                continue

            candidates: list[str] = []

            if isinstance(item, str):
                candidates.append(item)
            elif isinstance(item, dict):
                for key in ("file_path", "path", "file", "text", "source"):
                    val = item.get(key)
                    if val and isinstance(val, str):
                        candidates.append(val)
            else:
                # ANY object — check .text FIRST (this is where KB puts the path)
                text_val = getattr(item, "text", None)
                if text_val and isinstance(text_val, str) and text_val.strip():
                    candidates.append(text_val.strip())

                # Check .data dict
                data_dict = getattr(item, "data", None)
                if data_dict and isinstance(data_dict, dict):
                    for key in ("file_path", "path", "file", "text", "source"):
                        val = data_dict.get(key)
                        if val and isinstance(val, str) and val.strip():
                            candidates.append(val.strip())

                # Check .path attribute
                path_val = getattr(item, "path", None)
                if path_val is not None:
                    path_str = str(path_val).strip()
                    if path_str and path_str != "None":
                        candidates.append(path_str)

            # Validate all candidates
            for candidate in candidates:
                for line in candidate.split("\n"):
                    line = line.strip()
                    if not line or line == "None":
                        continue
                    try:
                        p = Path(line)
                        if p.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                            if p.exists():
                                if p not in paths:
                                    paths.append(p)
                            else:
                                logger.warning(f"[OCR] Path does not exist: {p}")
                    except (OSError, ValueError) as e:
                        logger.warning(f"[OCR] Invalid path '{line}': {e}")

        return paths

    # ══════════════════════════════════════════════════════════
    #  FILE DISPATCH
    # ══════════════════════════════════════════════════════════

    def _extract_file(self, path: Path) -> list[Data]:
        ext = path.suffix.lower()
        dispatch = {
            ".pdf": self._extract_pdf, ".docx": self._extract_docx,
            ".pptx": self._extract_pptx, ".xlsx": self._extract_xlsx,
            ".csv": self._extract_csv, ".txt": self._extract_txt,
        }
        if ext in dispatch:
            return dispatch[ext](path)
        elif ext in self.IMAGE_EXTENSIONS:
            return self._extract_image(path)
        return []

    def _extract_pdf(self, path: Path) -> list[Data]:
        import fitz  # PyMuPDF

        docs = []
        page_errors = []
        pdf_doc = fitz.open(str(path))
        total_pages = len(pdf_doc)

        for page_num in range(total_pages):
            page = pdf_doc[page_num]
            display_page = page_num + 1

            # Step 1: try native text extraction
            native_text = (page.get_text("text") or "").strip()
            if len(native_text) >= self.min_native_text_length:
                docs.append(Data(text=native_text, data={
                    "source_file": path.name, "file_path": str(path),
                    "page_number": display_page, "total_pages": total_pages,
                    "file_type": "pdf", "extraction_method": "native",
                }))
                continue

            # Step 2: scanned page — render to image and OCR via connected LLM
            if not self.llm:
                page_errors.append(f"Page {display_page}: scanned page but no LLM connected for OCR")
                continue

            try:
                ocr_text = self._ocr_pdf_page(pdf_doc, page_num)
                if ocr_text:
                    docs.append(Data(text=ocr_text, data={
                        "source_file": path.name, "file_path": str(path),
                        "page_number": display_page, "total_pages": total_pages,
                        "file_type": "pdf", "extraction_method": "llm_vision",
                    }))
                else:
                    page_errors.append(f"Page {display_page}: LLM returned empty text")
            except Exception as e:
                page_errors.append(f"Page {display_page}: {type(e).__name__}: {e}")

        pdf_doc.close()

        if not docs and page_errors:
            raise RuntimeError(
                f"All {total_pages} page(s) failed extraction. "
                f"Errors: {'; '.join(page_errors[:5])}"
            )

        return docs

    def _ocr_pdf_page(self, pdf_doc, page_index: int) -> str | None:
        """Render a PDF page to PNG and send to the connected LLM for OCR."""
        import fitz

        page = pdf_doc[page_index]
        zoom = self.ocr_dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        png_bytes = pix.tobytes("png")
        return self._llm_vision_extract(png_bytes, "image/png")

    def _extract_image(self, path: Path) -> list[Data]:
        if not self.llm:
            self.log("No LLM connected — cannot OCR images.")
            return []

        mime, _ = mimetypes.guess_type(str(path))
        text = self._llm_vision_extract(path.read_bytes(), mime or "image/png")
        if text:
            return [Data(text=text, data={
                "source_file": path.name, "file_path": str(path),
                "page_number": 1, "file_type": "image", "extraction_method": "llm_vision",
            })]
        return []

    def _extract_docx(self, path: Path) -> list[Data]:
        from docx import Document as DocxDocument
        doc = DocxDocument(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        if self.extract_tables:
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        parts.append(row_text)
        full_text = "\n".join(parts)
        if full_text.strip():
            return [Data(text=full_text, data={
                "source_file": path.name, "file_path": str(path),
                "page_number": 1, "file_type": "docx", "extraction_method": "native",
            })]
        return []

    def _extract_pptx(self, path: Path) -> list[Data]:
        from pptx import Presentation
        prs = Presentation(str(path))
        docs = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
                if self.extract_tables and shape.has_table:
                    for row in shape.table.rows:
                        rt = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                        if rt:
                            texts.append(rt)
            if texts:
                slide_text = "\n".join(texts)
                docs.append(Data(text=slide_text, data={
                    "source_file": path.name, "file_path": str(path),
                    "page_number": slide_num, "total_pages": len(prs.slides),
                    "file_type": "pptx", "extraction_method": "native",
                }))
        return docs

    def _extract_xlsx(self, path: Path) -> list[Data]:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), data_only=True)
        docs = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) if c is not None else "" for c in row]
                rt = " | ".join(v for v in vals if v)
                if rt:
                    rows.append(rt)
            if rows:
                docs.append(Data(text="\n".join(rows), data={
                    "source_file": path.name, "file_path": str(path),
                    "sheet_name": sheet_name, "page_number": 1,
                    "file_type": "xlsx", "extraction_method": "native",
                }))
        return docs

    def _extract_csv(self, path: Path) -> list[Data]:
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
        if text:
            return [Data(text=text, data={
                "source_file": path.name, "file_path": str(path),
                "page_number": 1, "file_type": "csv", "extraction_method": "native",
            })]
        return []

    def _extract_txt(self, path: Path) -> list[Data]:
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
        if text:
            return [Data(text=text, data={
                "source_file": path.name, "file_path": str(path),
                "page_number": 1, "file_type": "txt", "extraction_method": "native",
            })]
        return []

    # ══════════════════════════════════════════════════════════
    #  LLM VISION OCR (model-agnostic via LangChain)
    # ══════════════════════════════════════════════════════════

    def _llm_vision_extract(self, image_bytes: bytes, mime_type: str) -> str | None:
        """Send image to the connected LLM via LangChain multimodal messages."""
        from langchain_core.messages import HumanMessage

        b64_data = base64.b64encode(image_bytes).decode("utf-8")

        message = HumanMessage(content=[
            {
                "type": "text",
                "text": (
                    "Extract ALL text from this image accurately. "
                    "Preserve the original structure: headings, paragraphs, tables, lists. "
                    "Format tables as markdown tables. "
                    "Return ONLY the extracted text, no commentary or explanation."
                ),
            },
            {
                "type": "image_url",
                "image_url": {"url": f"data:{mime_type};base64,{b64_data}"},
            },
        ])

        response = self.llm.invoke([message])
        raw = response.content if hasattr(response, "content") else str(response)
        text = raw.strip()
        return text if text else None
