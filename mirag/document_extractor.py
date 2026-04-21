"""Document text extraction service.

Extracts text content from 22+ file types for document Q&A.
Ported from MiBuddy's extract_text_from_file() with agentcore storage integration.

PDF handling is a per-page hybrid: PyMuPDF native text extraction first;
pages with < 50 chars of native text are routed through Azure Document
Intelligence (`prebuilt-read` model) via the companion `ocr_extractor.py`
helper. A PDF with mix of text and scanned pages has both captured and
merged in page order.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path

from loguru import logger

# Per-page threshold for deciding a page is scanned. Matches
# components/OCR/ocr_extractor.py's min_native_text_length default.
_PDF_MIN_NATIVE_CHARS_PER_PAGE = 50

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}

SUPPORTED_DOC_EXTENSIONS = {
    # Documents
    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".xlsm",
    # Text / Markup
    ".txt", ".md", ".csv", ".html", ".tex",
    # Code
    ".py", ".js", ".ts", ".java", ".cpp", ".c", ".cs", ".go",
    ".rb", ".php", ".sh", ".css", ".json",
}


async def read_file_bytes(file_path: str) -> bytes:
    """Read file bytes from storage.

    Tries multiple storage locations:
    1. Main storage (files uploaded via /api/files)
    2. MiBuddy dedicated container (uploads folder)

    file_path format: "{user_id}/{filename}".
    """
    # Try MiBuddy container first (with full path)
    try:
        from agentcore.services.mibuddy.docqa_storage import get_file_by_path
        data = await get_file_by_path(file_path)
        logger.info(f"[DocExtractor] Read {len(data)} bytes from MiBuddy container: {file_path}")
        return data
    except Exception as e:
        logger.debug(f"[DocExtractor] Not in MiBuddy container ({e}), trying main storage: {file_path}")

    # Fallback: try main storage
    parts = file_path.replace("\\", "/").split("/", 1)
    if len(parts) == 2:
        agent_id, file_name = parts
    else:
        agent_id, file_name = "", file_path

    try:
        from agentcore.services.deps import get_storage_service
        storage = get_storage_service()
        return await storage.get_file(agent_id=agent_id, file_name=file_name)
    except Exception:
        raise FileNotFoundError(f"File not found in any storage: {file_path}")


def extract_text_from_bytes(file_bytes: bytes, file_ext: str) -> str:
    """Extract text content from file bytes based on extension.

    Args:
        file_bytes: Raw file content.
        file_ext: File extension including dot (e.g. ".pdf").

    Returns:
        Extracted text content.
    """
    file_ext = file_ext.lower()

    try:
        # --- Text / Code files ---
        if file_ext in {
            ".txt", ".md", ".csv", ".c", ".cpp", ".cs", ".css", ".go",
            ".java", ".js", ".json", ".php", ".py", ".rb", ".sh",
            ".tex", ".ts", ".html",
        }:
            text = file_bytes.decode("utf-8", errors="ignore")
            if file_ext == ".html":
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(text, "html.parser")
                text = soup.get_text(separator="\n")
            return text

        # --- PDF ---
        # Note: PDFs are handled in extract_text() (async) so we can OCR
        # scanned pages. This sync path is kept only as a safety fallback if
        # the caller bypasses extract_text().
        if file_ext == ".pdf":
            return _extract_pdf_native_only(file_bytes)

        # --- Word (DOCX) ---
        if file_ext == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
            # Tables — joined cell-by-cell per row with " | " separator so the
            # LLM can tell columns apart when it sees them in the prompt.
            # Matches the pattern in components/OCR/ocr_extractor.py.
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        c.text.strip() for c in row.cells if c.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)
            return "\n".join(parts)

        # --- PowerPoint (PPTX) ---
        if file_ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in paragraph.runs)
                            if line.strip():
                                text += line + "\n"
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text += cell.text + "\n"
            return text

        # --- Excel (XLS, XLSX, XLSM) ---
        if file_ext in {".xls", ".xlsx", ".xlsm"}:
            import pandas as pd
            file_stream = io.BytesIO(file_bytes)
            engine = "xlrd" if file_ext == ".xls" else "openpyxl"
            try:
                xls = pd.ExcelFile(file_stream, engine=engine)
            except Exception:
                file_stream.seek(0)
                xls = pd.ExcelFile(file_stream)

            parts: list[str] = []
            for sheet_idx, sheet_name in enumerate(xls.sheet_names, start=1):
                try:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    parts.append(f"\n--- Sheet {sheet_idx}: {sheet_name} ---\n")
                    parts.append(df.to_csv(index=False))
                except Exception as e:
                    logger.warning(f"Skipping Excel sheet '{sheet_name}': {e}")
            return "\n".join(parts).strip()

        return f"[Unsupported file type: {file_ext}]"

    except Exception as e:
        logger.error(f"Error extracting text from {file_ext}: {e}")
        return f"[ERROR] Unable to process file ({file_ext}). Reason: {str(e)}"


def _extract_pdf_native_only(file_bytes: bytes) -> str:
    """PyMuPDF native text extraction — no OCR. Used by the sync path.

    Returns an empty string if every page is scanned — the async path is what
    handles scanned PDFs. PyPDF2 is retained as a fallback if PyMuPDF fails
    to open the file (e.g. corrupt / unusual format).
    """
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        parts: list[str] = []
        for page in doc:
            # fitz.get_text("text") returns str; the `str()` cast pins the type
            # for type-checkers that see the union stub.
            page_text = str(page.get_text("text") or "")
            if page_text.strip():
                parts.append(page_text)
        doc.close()
        return "\n\n".join(parts)
    except Exception as exc:  # noqa: BLE001
        logger.warning(f"[DocExtractor] PyMuPDF failed, falling back to PyPDF2: {exc}")
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(file_bytes))
            parts = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    parts.append(page_text)
            return "\n\n".join(parts)
        except Exception as exc2:  # noqa: BLE001
            logger.error(f"[DocExtractor] Both PyMuPDF and PyPDF2 failed: {exc2}")
            return ""


async def _extract_pdf_with_ocr(file_bytes: bytes, display_name: str) -> str:
    """Per-page hybrid PDF extraction: PyMuPDF native + ADI OCR for scanned pages.

    - Pages with ≥ _PDF_MIN_NATIVE_CHARS_PER_PAGE non-whitespace chars use native text.
    - Remaining pages (empty or nearly-empty) are sent to Azure Document
      Intelligence in a SINGLE call scoped with `pages=` so we only pay for
      the scanned pages.
    - Native + OCR pages are merged back in page order.
    """
    from agentcore.services.mibuddy.ocr_extractor import ocr_pdf_pages, is_ocr_configured

    try:
        import fitz  # PyMuPDF
    except Exception as exc:  # noqa: BLE001
        logger.error(f"[DocExtractor] PyMuPDF unavailable: {exc}")
        return _extract_pdf_native_only(file_bytes)

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:  # noqa: BLE001
        logger.warning(f"[DocExtractor] PyMuPDF could not open {display_name}: {exc}")
        return _extract_pdf_native_only(file_bytes)

    native_by_page: dict[int, str] = {}
    scanned_pages: list[int] = []
    total_pages = doc.page_count

    for page_index, page in enumerate(doc, start=1):  # 1-based per ADI convention
        page_text = str(page.get_text("text") or "")
        if len(page_text.strip()) >= _PDF_MIN_NATIVE_CHARS_PER_PAGE:
            native_by_page[page_index] = page_text
        else:
            scanned_pages.append(page_index)

    doc.close()

    ocr_by_page: dict[int, str] = {}
    if scanned_pages:
        if is_ocr_configured():
            logger.info(
                f"[DocExtractor] {display_name}: "
                f"{len(scanned_pages)} scanned page(s) {scanned_pages} → ADI"
            )
            ocr_by_page = await ocr_pdf_pages(file_bytes, scanned_pages)
            missing = [p for p in scanned_pages if p not in ocr_by_page]
            if missing:
                logger.warning(
                    f"[DocExtractor] {display_name}: ADI returned no text for "
                    f"{len(missing)} page(s): {missing}"
                )
        else:
            logger.warning(
                f"[DocExtractor] {display_name}: "
                f"{len(scanned_pages)} scanned page(s) skipped — ADI not configured"
            )
    else:
        logger.info(
            f"[DocExtractor] {display_name}: 0 scanned pages, skipping ADI "
            f"({len(native_by_page)}/{total_pages} native)"
        )

    final_parts: list[str] = []
    for page_index in range(1, total_pages + 1):
        if page_index in native_by_page:
            final_parts.append(native_by_page[page_index])
        elif page_index in ocr_by_page:
            final_parts.append(ocr_by_page[page_index])
        # Genuinely empty or OCR-failed pages are skipped silently.

    return "\n\n".join(final_parts)


async def extract_text(file_path: str) -> str:
    """Read a file from storage and extract its text content.

    Args:
        file_path: Storage-relative path (e.g. "{user_id}/{filename}").

    Returns:
        Extracted text content.
    """
    ext = Path(file_path).suffix.lower()

    if ext in IMAGE_EXTENSIONS:
        return f"[Image file: {Path(file_path).name}]"

    if ext not in SUPPORTED_DOC_EXTENSIONS:
        return f"[Unsupported file type: {ext}]"

    try:
        file_bytes = await read_file_bytes(file_path)
        if ext == ".pdf":
            # Async path: per-page hybrid native + OCR.
            text = await _extract_pdf_with_ocr(file_bytes, Path(file_path).name)
        else:
            text = extract_text_from_bytes(file_bytes, ext)
            # OCR fallback for DOCX / PPTX: if python-docx / python-pptx
            # returned nothing, the file likely contains only embedded
            # images (scanned page pasted into Word, rasterized-slide
            # Keynote export, etc.). ADI natively accepts both formats.
            if not text.strip() and ext in {".docx", ".pptx"}:
                text = await _ocr_office_fallback(
                    file_bytes, Path(file_path).name, ext,
                )
        if not text.strip():
            return f"[No text content extracted from: {Path(file_path).name}]"
        return text
    except Exception as e:
        logger.error(f"Failed to extract text from {file_path}: {e}")
        return f"[ERROR] Failed to read file: {str(e)}"


async def _ocr_office_fallback(file_bytes: bytes, display_name: str, ext: str) -> str:
    """Whole-file ADI fallback for DOCX / PPTX (match MiBuddy's wider OCR net).

    Triggered only when the native parser (`python-docx` / `python-pptx`)
    returned empty — typically because the document's content is one or
    more embedded images rather than native text/runs.
    """
    from agentcore.services.mibuddy.ocr_extractor import (
        is_ocr_configured,
        ocr_office_bytes,
    )

    if not is_ocr_configured():
        logger.warning(
            f"[DocExtractor] {display_name}: native {ext} extractor returned "
            f"empty, ADI not configured — file will appear empty"
        )
        return ""

    logger.info(
        f"[DocExtractor] {display_name}: native {ext} extractor returned "
        f"empty → falling back to ADI (embedded-image content)"
    )
    return await ocr_office_bytes(file_bytes)
